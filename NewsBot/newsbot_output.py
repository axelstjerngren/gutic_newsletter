from pptx import Presentation
from pptx.util import Cm, Pt
import os
from pptx_library import python_pptx_funcs



class pptx_output():
	def __init__(self, titles, articles, pics, savefolder):

		self.prs = Presentation(os.getcwd() + "/resources/rough_template.pptx")
		for slide in self.prs.slides:
			self.slide = slide
		self.add_titles(titles)
		self.add_articles(articles)
		self.add_pictures(pics)
		self.save_presentation(savefolder)

	def add_titles(self, titles):
		titlelist = ["TITLE1","TITLE2","TITLE3","TITLE4","TITLE5"]
		for title, target in zip(titles, titlelist):
			python_pptx_funcs.run_text(self.slide,title,target)

	def add_articles(self,articles):
		textlist = ["TEXT1","TEXT2","TEXT3","TEXT4","TEXT5"]
		for text, target in zip(articles, textlist):
			python_pptx_funcs.run_text(self.slide,text,target)

	def add_pictures(self,pics):
		counter = 0
		if len(pics) == 0:
			return
		for shape in self.slide.shapes:
			if counter == 5:
				break
			if shape.is_placeholder:
				phf = shape.placeholder_format
				print(phf.type)
				if "PICTURE" in str(phf.type):
					image = pics[counter]
					picture = shape.insert_picture(image)
					counter += 1


	def save_presentation(self, savefolder):
		self.prs.save(savefolder + "/NewsLetter.pptx")
		print("Presentation saved at: "+ savefolder)




